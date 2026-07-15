from __future__ import annotations

import re
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable

QUESTION_HEADINGS = re.compile(r"(model|tutorial|assignment|quiz|previous.?year|end.?of.?chapter).*questions?|question bank|exercises?", re.I)
HEADING = re.compile(r"^(?:module|unit|chapter|topic|[0-9]+(?:\.[0-9]+)*[.:])\s*.+", re.I)


@dataclass
class Chunk:
    chunk_id: str
    text: str
    source: str
    module: str = ""
    chapter: str = ""
    topic: str = ""
    page_number: int = 0
    corpus: str = "content"

    def metadata(self) -> dict:
        return asdict(self)


def _read_pages(path: Path) -> list[str]:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return [path.read_text(encoding="utf-8", errors="ignore")]
    if suffix == ".pdf":
        from pypdf import PdfReader
        return [(page.extract_text() or "") for page in PdfReader(path).pages]
    if suffix == ".docx":
        from docx import Document
        return ["\n".join(p.text for p in Document(path).paragraphs)]
    if suffix == ".pptx":
        from pptx import Presentation
        return ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in Presentation(path).slides]
    raise ValueError(f"Unsupported file type: {path.suffix}")


def _split(text: str, size: int, overlap: int) -> Iterable[str]:
    paragraphs = [re.sub(r"\s+", " ", p).strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > size:
            for sentence in re.split(r"(?<=[.!?])\s+", paragraph):
                if current and len(current) + len(sentence) + 1 > size:
                    yield current
                    current = current[-overlap:] + " " + sentence if overlap else sentence
                else:
                    current = f"{current} {sentence}".strip()
            continue
        if current and len(current) + len(paragraph) + 1 > size:
            yield current
            current = current[-overlap:] + "\n" + paragraph if overlap else paragraph
        else:
            current = f"{current}\n{paragraph}".strip()
    if current:
        yield current


def _heading(label: str, is_question_heading: bool) -> bool:
    return is_question_heading or bool(HEADING.match(label)) or (len(label) < 90 and label.isupper())


def _question_parts(text: str) -> list[str]:
    parts = re.split(r"(?=(?:^|\s)\d+\.\s+)", re.sub(r"\s+", " ", text).strip())
    return [part.strip() for part in parts if re.match(r"^\d+\.\s+", part.strip())]


def ingest_file(path: str | Path, chunk_size: int = 900, chunk_overlap: int = 150) -> list[Chunk]:
    path = Path(path)
    module = chapter = topic = ""
    corpus = "content"
    chunks: list[Chunk] = []
    def save(text: str, page_number: int) -> None:
        text = "\n".join(text.splitlines()).strip()
        if not text:
            return
        parts = _question_parts(text) if corpus == "questions" else list(_split(text, chunk_size, chunk_overlap))
        for part in parts:
            chunks.append(Chunk(f"{path.stem}-{page_number}-{len(chunks)}", part, path.name, module, chapter, topic, page_number, corpus))

    for page_number, page in enumerate(_read_pages(path), 1):
        lines = [line.strip() for line in page.splitlines() if line.strip()]
        body: list[str] = []
        for line in lines:
            label = re.sub(r"\s+", " ", line)
            is_question_heading = bool(QUESTION_HEADINGS.search(label))
            if _heading(label, is_question_heading):
                save("\n".join(body), page_number)
                body = []
                if is_question_heading:
                    corpus = "questions"
                elif "module" in label.lower() or "unit" in label.lower():
                    corpus = "content"
                if "module" in label.lower() or "unit" in label.lower(): module = label
                elif "chapter" in label.lower(): chapter = label
                else: topic = label
            body.append(label)
        save("\n".join(body), page_number)
    return chunks
